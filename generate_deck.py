import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Premium SaaS Design System Colors (Inspired by Tailwind CSS & Github Dark)
    bg_color = RGBColor(10, 12, 16)          # Deep cosmic black (#0A0C10)
    card_bg = RGBColor(18, 22, 32)           # Dark slate gray-blue (#121620)
    card_border = RGBColor(33, 41, 56)       # Subtle slate border (#212938)
    
    accent_purple = RGBColor(139, 92, 246)   # Radiant Violet (#8B5CF6)
    accent_red = RGBColor(234, 88, 12)       # Safety Orange-Red (#EA580C)
    accent_green = RGBColor(16, 185, 129)    # Emerald Green (#10B981)
    
    text_offwhite = RGBColor(245, 247, 250)  # Pure soft off-white (#F5F7FA)
    
    # Muted opacities simulated against #0A0C10
    text_white_70 = RGBColor(170, 175, 185)  # 70% opacity
    text_white_50 = RGBColor(120, 125, 135)  # 50% opacity
    text_white_30 = RGBColor(75, 80, 90)     # 30% opacity
    text_white_15 = RGBColor(38, 40, 46)     # 15% opacity
    
    accent_purple_25 = RGBColor(41, 31, 69)  # 25% opacity Violet
    accent_purple_15 = RGBColor(29, 23, 48)  # 15% opacity Violet
    accent_red_20 = RGBColor(56, 26, 21)     # 20% opacity Orange-Red
    
    # Typography
    font_header = "Georgia"
    font_body = "Arial"  # Grotesk representation
    font_mono = "Consolas"
    
    # Helper to set full-bleed background shape covering the entire slide (guarantees solid theme color)
    def set_dark_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background() # No border line

    # Helper to add a standardized header
    def add_slide_header(slide, title_text, subtitle_text=""):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.3))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Title paragraph (Georgia)
        p1 = tf.paragraphs[0]
        p1.text = title_text
        p1.font.name = font_header
        p1.font.size = Pt(32)
        p1.font.bold = False
        p1.font.color.rgb = text_offwhite
        p1.space_after = Pt(4)
        
        # Subtitle paragraph (Grotesk)
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = font_body
            p2.font.size = Pt(14.5)
            p2.font.color.rgb = text_white_70
            
    # Helper to draw a styled card
    def add_card(slide, left, top, width, height, bg_rgb=card_bg, border_rgb=card_border):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_rgb
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1)
        if shape.adjustments:
            try:
                shape.adjustments[0] = 0.04
            except:
                pass
        return shape

    # Helper to add pill badges
    def add_pill_badge(slide, left, top, width, height, text, bg_rgb, text_rgb):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_rgb
        shape.line.fill.background() # No border
        if shape.adjustments:
            try:
                shape.adjustments[0] = 0.5
            except:
                pass
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.alignment = PP_ALIGN.CENTER
        p.font.name = font_body
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = text_rgb

    # =========================================================================
    # SLIDE 1: Title Slide (Clinical research style)
    # =========================================================================
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide1)
    
    # 4-5 nodes connected by thin #8B5CF6 lines at 30% opacity, upper-right, bleeding off-canvas
    nodes = [
        (Inches(10.2), Inches(0.8)),
        (Inches(11.5), Inches(2.0)),
        (Inches(9.2), Inches(3.0)),
        (Inches(10.8), Inches(4.0)),
        (Inches(12.5), Inches(4.5))
    ]
    connections = [(0, 1), (1, 2), (2, 3), (3, 4), (1, 4)]
    
    # Draw connections (30% opacity purple)
    for start_idx, end_idx in connections:
        n1 = nodes[start_idx]
        n2 = nodes[end_idx]
        line = slide1.shapes.add_connector(
            1,
            n1[0] + Pt(4), n1[1] + Pt(4),
            n2[0] + Pt(4), n2[1] + Pt(4)
        )
        line.line.color.rgb = accent_purple_25
        line.line.width = Pt(1)
            
    # Draw nodes
    for i, (x, y) in enumerate(nodes):
        node = slide1.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.12), Inches(0.12))
        node.fill.solid()
        if i == 1:
            node.fill.fore_color.rgb = accent_purple
            node.line.color.rgb = text_offwhite
            node.line.width = Pt(1)
        else:
            node.fill.fore_color.rgb = text_white_30
            node.line.color.rgb = bg_color
            node.line.width = Pt(1)

    # Title & Subtitle in lower third (Separated to prevent line spacing overlaps)
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p_title = tf.paragraphs[0]
    p_title.text = "MindGraph"
    p_title.font.name = font_header
    p_title.font.size = Pt(88)
    p_title.font.bold = False
    p_title.font.color.rgb = text_offwhite
    
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.3), Inches(1.0))
    stf = sub_box.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = 0
    p_sub = stf.paragraphs[0]
    p_sub.text = "Nobody burns out overnight. It's written in your data for weeks first."
    p_sub.font.name = font_header
    p_sub.font.size = Pt(22)
    p_sub.font.italic = True
    p_sub.font.color.rgb = text_white_70

    # Thin horizontal rule in #EA580C at 20% opacity above footer
    rule = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.3), Inches(11.3), Pt(1))
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent_red_20
    rule.line.fill.background()

    # Footer strip at the very bottom-left (Grotesk, small caps)
    credits_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5))
    ctf = credits_box.text_frame
    ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]
    cp.text = "DARK PULSE   ·   HACKHAZARDS '26   ·   HUMAN EXPERIENCE & PRODUCTIVITY TRACK\nMANOJ H.G (LEAD)   ·   DILEEP MK   ·   SNEHA   ·   CHINMAY J C"
    cp.font.name = font_body
    cp.font.size = Pt(9.5)
    cp.font.color.rgb = text_white_30

    # =========================================================================
    # SLIDE 2: The Problem
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "We stopped noticing ourselves a long time ago.", "The signs were there. We just built apps that don't talk to each other.")
    
    # Left Column: "77%" Number + Description (aligned to avoid overlaps)
    num_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(6.5), Inches(2.0))
    ntf = num_box.text_frame
    ntf.word_wrap = True
    ntf.margin_left = ntf.margin_top = ntf.margin_right = ntf.margin_bottom = 0
    np = ntf.paragraphs[0]
    np.text = "77%"
    np.font.name = font_header
    np.font.size = Pt(130)
    np.font.bold = False
    np.font.color.rgb = accent_red
    
    lbl_box = slide2.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(6.0), Inches(0.4))
    ltf = lbl_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = 0
    lp = ltf.paragraphs[0]
    lp.text = "REPORT BURNOUT SYMPTOMS"
    lp.font.name = font_body
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = text_white_50
    
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(6.5), Inches(2.0))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = dtf.margin_top = dtf.margin_right = dtf.margin_bottom = 0
    dp = dtf.paragraphs[0]
    dp.text = "Your sleep tracker knows you slept five hours. Your fitness app knows you didn't move. Your calendar knows six back-to-back calls. None of them know that together, that's the third week running."
    dp.font.name = font_body
    dp.font.size = Pt(13.5)
    dp.font.color.rgb = text_white_70

    # Right Column: Premium styled app panels
    right_x = Inches(8.3)
    apps = ["Sleep Tracker", "Movement Log", "Calendar App"]
    
    for idx, app_name in enumerate(apps):
        y_pos = Inches(2.2) + idx * Inches(1.3)
        # Add Card
        add_card(slide2, right_x, y_pos, Inches(4.2), Inches(0.55))
        
        # Overlay Pill text inside card
        bar_box = slide2.shapes.add_textbox(right_x + Inches(0.2), y_pos + Inches(0.12), Inches(3.8), Inches(0.3))
        btf = bar_box.text_frame
        btf.margin_left = btf.margin_top = btf.margin_right = btf.margin_bottom = 0
        bp = btf.paragraphs[0]
        bp.text = app_name.upper()
        bp.font.name = font_body
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = text_offwhite
        
        # Center connector (Precise rect representation)
        if idx < 2:
            line_y = y_pos + Inches(0.55)
            l = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x + Inches(2.1) - Pt(0.75), line_y, Pt(1.5), Inches(0.75))
            l.fill.solid()
            l.fill.fore_color.rgb = accent_red
            l.line.fill.background()

    # Closing line: Georgia italic
    closer_box = slide2.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
    cltf = closer_box.text_frame
    cltf.margin_left = cltf.margin_top = cltf.margin_right = cltf.margin_bottom = 0
    clp = cltf.paragraphs[0]
    clp.text = "The data to predict this already exists. It's just scattered across four apps that don't speak."
    clp.alignment = PP_ALIGN.CENTER
    clp.font.name = font_header
    clp.font.size = Pt(18)
    clp.font.italic = True
    clp.font.color.rgb = text_white_70

    # =========================================================================
    # SLIDE 3: The Solution
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "What if your habits could tell on each other?", "MindGraph treats your life like it actually is — one connected system, not four separate apps.")
    
    # Continuous connector line behind beats
    line_spine = slide3.shapes.add_connector(
        1,
        Inches(2.6), Inches(3.5),
        Inches(10.6), Inches(3.5)
    )
    line_spine.line.color.rgb = accent_purple
    line_spine.line.width = Pt(1.5)
    
    beat_w = Inches(3.6)
    beat_h = Inches(3.4)
    y_pos = Inches(2.3)
    
    # Beat 1 (Left) - styled card with Pill badge
    add_card(slide3, Inches(0.8), y_pos, beat_w, beat_h, border_rgb=card_border)
    add_pill_badge(slide3, Inches(1.0), y_pos + Inches(0.2), Inches(0.8), Inches(0.25), "01 / IN", accent_purple_25, accent_purple)
    
    box1 = slide3.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.6), beat_w - Inches(0.4), beat_h - Inches(0.8))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    p = tf1.paragraphs[0]
    p.text = "30 seconds"
    p.font.name = font_body
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_offwhite
    p.space_after = Pt(10)
    p = tf1.add_paragraph()
    p.text = "A 30-second daily log: sleep, mood, energy. Offline-first client cache ensures entries are preserved."
    p.font.name = font_body
    p.font.size = Pt(12.5)
    p.font.color.rgb = text_white_70

    # Beat 2 (Center) - Node graph card
    add_card(slide3, Inches(4.8), y_pos, beat_w, beat_h, border_rgb=accent_purple)
    add_pill_badge(slide3, Inches(5.0), y_pos + Inches(0.2), Inches(1.1), Inches(0.25), "02 / GRAPH", accent_purple_25, accent_purple)
    
    box2 = slide3.shapes.add_textbox(Inches(5.0), y_pos + Inches(0.6), beat_w - Inches(0.4), beat_h - Inches(0.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p = tf2.paragraphs[0]
    p.text = "Connected Node Graph"
    p.font.name = font_body
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_offwhite
    p.space_after = Pt(10)
    p = tf2.add_paragraph()
    p.text = "Every entry becomes a node. Not stored — related, to yesterday, to your patterns."
    p.font.name = font_body
    p.font.size = Pt(12.5)
    p.font.color.rgb = text_white_70
    
    # Graph Visual inside Center Card
    gy = y_pos + Inches(2.6)
    g_nodes = [(Inches(5.3), gy), (Inches(6.1), gy + Inches(0.25)), (Inches(6.8), gy - Inches(0.1)), (Inches(7.4), gy + Inches(0.2)), (Inches(8.0), gy - Inches(0.05))]
    g_conns = [(0, 1), (1, 2), (2, 3), (3, 4), (1, 3)]
    for s_i, e_i in g_conns:
        n1 = g_nodes[s_i]
        n2 = g_nodes[e_i]
        l = slide3.shapes.add_connector(1, n1[0]+Pt(3), n1[1]+Pt(3), n2[0]+Pt(3), n2[1]+Pt(3))
        l.line.color.rgb = accent_purple
        l.line.width = Pt(1)
    for x, y in g_nodes:
        n = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.1), Inches(0.1))
        n.fill.solid()
        n.fill.fore_color.rgb = accent_purple
        n.line.color.rgb = text_offwhite
        n.line.width = Pt(0.75)

    # Beat 3 (Right)
    add_card(slide3, Inches(8.8), y_pos, beat_w, beat_h, border_rgb=card_border)
    add_pill_badge(slide3, Inches(9.0), y_pos + Inches(0.2), Inches(0.9), Inches(0.25), "03 / OUT", accent_purple_25, accent_purple)
    
    box3 = slide3.shapes.add_textbox(Inches(9.0), y_pos + Inches(0.6), beat_w - Inches(0.4), beat_h - Inches(0.8))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
    p = tf3.paragraphs[0]
    p.text = "Burnout Risk — Live"
    p.font.name = font_body
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = text_offwhite
    p.space_after = Pt(10)
    p = tf3.add_paragraph()
    p.text = "Output: a live Burnout Risk Gauge, running BIS Score, and coaching notes that remember history."
    p.font.name = font_body
    p.font.size = Pt(12.5)
    p.font.color.rgb = text_white_70

    # =========================================================================
    # SLIDE 4: Tech Stack (Tectonic Stack cross-section)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "Nothing about this is a toy demo.", "The stack exists because relationships need a database built for relationships.")
    
    # Left third: Stack layout
    stack_x = Inches(0.8)
    stack_w = Inches(3.8)
    
    # Top Layer: Intelligence
    add_card(slide4, stack_x, Inches(2.0), stack_w, Inches(1.1), bg_rgb=accent_purple, border_rgb=accent_purple)
    box = slide4.shapes.add_textbox(stack_x + Inches(0.2), Inches(2.15), stack_w - Inches(0.4), Inches(0.8))
    bt = box.text_frame
    bp = bt.paragraphs[0]
    bp.text = "INTELLIGENCE LAYER"
    bp.font.name = font_mono
    bp.font.size = Pt(9)
    bp.font.bold = True
    bp.font.color.rgb = text_offwhite
    bp2 = bt.add_paragraph()
    bp2.text = "GPT-4 Wellness OS"
    bp2.font.name = font_body
    bp2.font.size = Pt(14)
    bp2.font.bold = True
    bp2.font.color.rgb = text_offwhite

    # Middle Layer: Engine
    add_card(slide4, stack_x, Inches(3.1), stack_w, Inches(1.4), bg_rgb=accent_purple_15, border_rgb=accent_purple)
    box = slide4.shapes.add_textbox(stack_x + Inches(0.2), Inches(3.25), stack_w - Inches(0.4), Inches(1.1))
    bt = box.text_frame
    bp = bt.paragraphs[0]
    bp.text = "GRAPH ENGINE"
    bp.font.name = font_mono
    bp.font.size = Pt(9)
    bp.font.bold = True
    bp.font.color.rgb = text_white_70
    bp2 = bt.add_paragraph()
    bp2.text = "Node.js + Neo4j AuraDB"
    bp2.font.name = font_body
    bp2.font.size = Pt(14)
    bp2.font.bold = True
    bp2.font.color.rgb = text_offwhite

    # Bottom Layer: Client
    add_card(slide4, stack_x, Inches(4.5), stack_w, Inches(1.8), bg_rgb=card_bg, border_rgb=card_border)
    box = slide4.shapes.add_textbox(stack_x + Inches(0.2), Inches(4.65), stack_w - Inches(0.4), Inches(1.6))
    bt = box.text_frame
    bp = bt.paragraphs[0]
    bp.text = "CLIENT INTERFACE"
    bp.font.name = font_mono
    bp.font.size = Pt(9)
    bp.font.bold = True
    bp.font.color.rgb = text_white_70
    bp2 = bt.add_paragraph()
    bp2.text = "Expo + React Native"
    bp2.font.name = font_body
    bp2.font.size = Pt(14)
    bp2.font.bold = True
    bp2.font.color.rgb = text_offwhite

    # Right annotations
    ann_x = Inches(5.0)
    ann_w = Inches(7.5)
    
    # Ann 1 (Intelligence)
    box = slide4.shapes.add_textbox(ann_x, Inches(2.0), ann_w, Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    ap = tf.paragraphs[0]
    ap.text = "Intelligence — GPT-4 reads the graph itself, not a single entry — a coaching note can reference three weeks of continuous data, not just today's mood."
    ap.font.name = font_body
    ap.font.size = Pt(13)
    ap.font.color.rgb = text_white_70
    apt = tf.add_paragraph()
    apt.text = "openai-api  ·  gpt-4  ·  gpt-3.5-turbo"
    apt.font.name = font_mono
    apt.font.size = Pt(9.5)
    apt.font.color.rgb = text_white_50
    apt.font.bold = True
    apt.space_before = Pt(3)

    # Ann 2 (Engine)
    box = slide4.shapes.add_textbox(ann_x, Inches(3.1), ann_w, Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    ap = tf.paragraphs[0]
    ap.text = "Engine — Node.js/Express backend writing directly into Neo4j AuraDB in Cypher. Habits aren't rows in a table. They're edges."
    ap.font.name = font_body
    ap.font.size = Pt(13)
    ap.font.color.rgb = text_white_70
    apt = tf.add_paragraph()
    apt.text = "neo4j-auradb  ·  cypher-query-language  ·  express.js"
    apt.font.name = font_mono
    apt.font.size = Pt(9.5)
    apt.font.color.rgb = text_white_50
    apt.font.bold = True
    apt.space_before = Pt(3)

    # Ann 3 (Client)
    box = slide4.shapes.add_textbox(ann_x, Inches(4.5), ann_w, Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    ap = tf.paragraphs[0]
    ap.text = "Client — Expo + React Native 0.85, TypeScript 6.0. Offline-first: log a bad day with no signal, sync later through AsyncStorage."
    ap.font.name = font_body
    ap.font.size = Pt(13)
    ap.font.color.rgb = text_white_70
    apt = tf.add_paragraph()
    apt.text = "expo-sdk-56  ·  react-native-0.85  ·  async-storage  ·  typescript-6"
    apt.font.name = font_mono
    apt.font.size = Pt(9.5)
    apt.font.color.rgb = text_white_50
    apt.font.bold = True
    apt.space_before = Pt(3)

    # Pointing lines
    connector_y_positions = [Inches(2.55), Inches(3.80), Inches(5.45)]
    for idx, cy in enumerate(connector_y_positions):
        conn = slide4.shapes.add_connector(1, stack_x + stack_w, cy, ann_x - Inches(0.15), cy)
        conn.line.color.rgb = text_white_30
        conn.line.width = Pt(1)

    # =========================================================================
    # SLIDE 5: Demo & Cypher Query
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide5)
    add_slide_header(slide5, "This is what a pattern looks like once you can see it.", "A real query, running against a real graph.")
    
    # Top 40%: Demo Capture box with accent border
    demo_x = Inches(0.8)
    demo_w = Inches(11.7)
    demo_h = Inches(2.2)
    demo_placeholder = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, demo_x, Inches(2.0), demo_w, demo_h)
    demo_placeholder.fill.solid()
    demo_placeholder.fill.fore_color.rgb = card_bg
    demo_placeholder.line.color.rgb = card_border
    demo_placeholder.line.width = Pt(1.5)
    
    tf_dp = demo_placeholder.text_frame
    tf_dp.word_wrap = True
    dpp = tf_dp.paragraphs[0]
    dpp.text = "DEMO CAPTURE"
    dpp.alignment = PP_ALIGN.CENTER
    dpp.font.name = font_body
    dpp.font.size = Pt(12)
    dpp.font.bold = True
    dpp.font.color.rgb = text_white_50
    
    # Bottom Left: Gauge Card
    left_card_x = Inches(0.8)
    card_y = Inches(4.4)
    card_w = Inches(5.6)
    card_h = Inches(2.0)
    add_card(slide5, left_card_x, card_y, card_w, card_h)
    add_pill_badge(slide5, left_card_x + Inches(0.25), card_y + Inches(0.2), Inches(1.2), Inches(0.25), "OS ENGINE", accent_purple_25, accent_purple)
    
    score_num_box = slide5.shapes.add_textbox(left_card_x + Inches(0.25), card_y + Inches(0.55), card_w - Inches(0.5), Inches(0.6))
    sntf = score_num_box.text_frame
    sntf.word_wrap = True
    sntf.margin_left = sntf.margin_top = sntf.margin_right = sntf.margin_bottom = 0
    p = sntf.paragraphs[0]
    p.text = "BIS Grade: B  (Score: 73%)"
    p.font.name = font_header
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = accent_green
    
    score_desc_box = slide5.shapes.add_textbox(left_card_x + Inches(0.25), card_y + Inches(1.15), card_w - Inches(0.5), Inches(0.8))
    sdtf = score_desc_box.text_frame
    sdtf.word_wrap = True
    sdtf.margin_left = sdtf.margin_top = sdtf.margin_right = sdtf.margin_bottom = 0
    p = sdtf.paragraphs[0]
    p.text = "Formula: 20% Mood + 20% Sleep (vs 8h) + 20% Consistency + 20% Productivity + 20% Stress. Recalculated live to warn before physical burnout."
    p.font.name = font_body
    p.font.size = Pt(11.5)
    p.font.color.rgb = text_white_70

    # Bottom Right: Node diagram
    right_card_x = Inches(6.9)
    add_card(slide5, right_card_x, card_y, card_w, card_h)
    add_pill_badge(slide5, right_card_x + Inches(0.25), card_y + Inches(0.2), Inches(1.4), Inches(0.25), "NEO4J PATHS", accent_purple_25, accent_purple)
    
    cypher_code_box = slide5.shapes.add_textbox(right_card_x + Inches(0.25), card_y + Inches(0.55), card_w - Inches(0.5), Inches(0.5))
    ccf = cypher_code_box.text_frame
    ccf.word_wrap = True
    ccf.margin_left = ccf.margin_top = ccf.margin_right = ccf.margin_bottom = 0
    p = ccf.paragraphs[0]
    p.text = "Sleep ➔ Mood ➔ Productivity  (+66 pts)\nStress ➔ Sleep Deprivation ➔ Burnout"
    p.font.name = font_mono
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = accent_purple
    
    cypher_desc_box = slide5.shapes.add_textbox(right_card_x + Inches(0.25), card_y + Inches(1.15), card_w - Inches(0.5), Inches(0.8))
    cdtf = cypher_desc_box.text_frame
    cdtf.word_wrap = True
    cdtf.margin_left = cdtf.margin_top = cdtf.margin_right = cdtf.margin_bottom = 0
    p = cdtf.paragraphs[0]
    p.text = "Three weeks of bad sleep doesn't sit in a table row. It traces a path straight to the score."
    p.font.name = font_body
    p.font.size = Pt(12)
    p.font.color.rgb = text_white_70

    # Monospace query footer
    cypher_text = slide5.shapes.add_textbox(right_card_x + Inches(0.25), card_y + Inches(1.65), card_w - Inches(0.5), Inches(0.3))
    ctf = cypher_text.text_frame
    ctf.word_wrap = True
    ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]
    cp.text = "MATCH (u:User)-[:SLEPT]->(s:Sleep) MATCH (u)-[:LOGGED]->(m:Mood {date: s.date})"
    cp.font.name = font_mono
    cp.font.size = Pt(8.5)
    cp.font.color.rgb = text_white_50

    # =========================================================================
    # SLIDE 6: Concluding Slide
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide6)
    add_slide_header(slide6, "Intercept burnout before it peaks.", "The twelve weeks of warning nobody currently gets.")
    
    # Left Side: spelled outer nodes layout
    cy = Inches(3.3)
    cx = Inches(2.3)
    
    g_nodes = [
        (Inches(0.8), Inches(1.8), "Sleep"),
        (Inches(3.9), Inches(1.8), "Mood"),
        (Inches(0.8), Inches(4.9), "Stress"),
        (Inches(3.9), Inches(4.9), "Habits")
    ]
    
    hub_cx = cx + Inches(0.6)
    hub_cy = cy + Inches(0.6)
    
    for x, y, label in g_nodes:
        out_cx = x + Inches(0.55)
        out_cy = y + Inches(0.55)
        conn = slide6.shapes.add_connector(1, hub_cx, hub_cy, out_cx, out_cy)
        conn.line.color.rgb = accent_purple_25
        conn.line.width = Pt(1.5)
        
    for x, y, label in g_nodes:
        nd = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.1), Inches(1.1))
        nd.fill.solid()
        nd.fill.fore_color.rgb = card_bg
        nd.line.color.rgb = accent_purple
        nd.line.width = Pt(1.2)
        
        nt = nd.text_frame
        nt.margin_left = nt.margin_top = nt.margin_right = nt.margin_bottom = 0
        nt.word_wrap = False
        np = nt.paragraphs[0]
        np.text = label
        np.alignment = PP_ALIGN.CENTER
        np.font.name = font_body
        np.font.size = Pt(11)
        np.font.bold = True
        np.font.color.rgb = text_white_70

    # Central Node (MindGraph Hub Logo style)
    center_hub = slide6.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(1.2), Inches(1.2))
    center_hub.fill.solid()
    center_hub.fill.fore_color.rgb = accent_purple
    center_hub.line.color.rgb = text_offwhite
    center_hub.line.width = Pt(1.5)
    
    htf = center_hub.text_frame
    htf.margin_left = htf.margin_top = htf.margin_right = htf.margin_bottom = 0
    hp = htf.paragraphs[0]
    hp.text = "Mind\nGraph"
    hp.alignment = PP_ALIGN.CENTER
    hp.font.name = font_header
    hp.font.size = Pt(13)
    hp.font.bold = True
    hp.font.color.rgb = text_offwhite

    # Right Side: Team & Contact Card
    card_x = Inches(6.0)
    card_y = Inches(2.1)
    card_w = Inches(6.5)
    card_h = Inches(4.4)
    add_card(slide6, card_x, card_y, card_w, card_h, border_rgb=accent_purple)
    
    box = slide6.shapes.add_textbox(card_x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "TEAM DARK PULSE"
    p.font.name = font_mono
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = accent_purple
    p.space_after = Pt(14)
    
    contacts = [
        ("Manoj H.G (Lead)", "manojhg321@gmail.com"),
        ("Dileep MK", "dileep.m.k126@gmail.com"),
        ("Sneha", "acharyasneha2005@gmail.com"),
        ("Chinmay J C", "chinmaychoudhari620@gmail.com")
    ]
    for name, email in contacts:
        p = tf.add_paragraph()
        p.text = f"{name}   ·   {email}"
        p.font.name = font_body
        p.font.size = Pt(13)
        p.font.color.rgb = text_offwhite
        p.space_after = Pt(8)
        
    p = tf.add_paragraph()
    p.text = "\nREPOSITORY: github.com/sneha-s2005/MindGraph"
    p.font.name = font_mono
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = text_white_50

    p = tf.add_paragraph()
    p.text = "LIVE DEPLOYMENT: mindgraph-app.onrender.com"
    p.font.name = font_mono
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = text_white_50

    # Bookend closing tagline
    closer_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
    cltf = closer_box.text_frame
    cltf.margin_left = cltf.margin_top = cltf.margin_right = cltf.margin_bottom = 0
    clp = cltf.paragraphs[0]
    clp.text = "Nobody burns out overnight. It's written in your data for weeks first."
    clp.alignment = PP_ALIGN.CENTER
    clp.font.name = font_header
    clp.font.size = Pt(18)
    clp.font.italic = True
    clp.font.color.rgb = text_white_70

    # Save the presentation
    output_filename = "MindGraph_Presentation_Final_v8.pptx"
    prs.save(output_filename)
    print(f"Final presentation successfully created: {output_filename}")

if __name__ == "__main__":
    create_presentation()
